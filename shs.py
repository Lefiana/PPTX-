import os
import copy
import io
import re
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm
from PIL import Image, ImageOps
import unicodedata

# ==========================================
# --- CONFIGURATION & PATHS ---
# ==========================================
EXCEL_PATH = "Graduating_SHS.xlsx"
MASTER_STUDENT_DIR = Path(r"G:\Grad\SHS")
TEMPLATE_PPTX = "Commencement Exercises - PPT Template.pptx.pptx"
OUTPUT_PPTX = "Final_Graduation_SHS.pptx"
TEMPLATE_SLIDE_INDEX = 0
PLACEHOLDER_IMAGE = "no_picture.png"

# ==========================================
# --- SHAPE TARGETS & POSITIONS ---
# ==========================================
TARGET_SURNAME_SHAPE = "TextBox 6"
TARGET_FIRSTNAME_SHAPE = "TextBox 7"
TARGET_COURSE_SHAPE = "TextBox 8"

# Based on PowerPoint 'Format Picture' coordinates
IMG_LEFT = Cm(29.79)   
IMG_TOP = Cm(2.58)     
IMG_WIDTH = Cm(17.27)  
IMG_HEIGHT = Cm(23.43) 

# ==========================================
# --- DATA PROCESSING ---
# ==========================================
def normalize_text(text):
    """Converts Ñ to N and removes special characters for matching."""
    # Convert Ñ to N
    text = unicodedata.normalize('NFKD', text).encode('ASCII', 'ignore').decode('utf-8')
    # Remove underscores and punctuation, convert to lowercase
    return re.sub(r'[^\w\s]', ' ', text).lower()
def parse_name_for_slide(full_name: str):
    """
    Splits the Excel name into Surname and Firstname safely.
    Ideal format in Excel: "De Guzman, Juan"
    """
    if ',' in full_name:
        parts = full_name.split(',', 1)
        return parts[0].strip().upper(), parts[1].strip().title()
    
    # Fallback if no comma is used (Smart PH Name detection)
    words = full_name.split()
    if len(words) == 1:
        return words[0].upper(), ""
        
    prefixes = ['de', 'del', 'de la', 'dela', 'san', 'santa', 'sto', 'sto.']
    
    if len(words) >= 2 and words[-2].lower() in prefixes:
        surname = " ".join(words[-2:])
        first = " ".join(words[:-2])
    elif len(words) >= 3 and " ".join(words[-3:-1]).lower() in prefixes:
        surname = " ".join(words[-3:])
        first = " ".join(words[:-3])
    else:
        surname = words[-1]
        first = " ".join(words[:-1])
        
    return surname.upper(), first.title()

def get_validated_students(excel_path, master_dir):
    print("--- PHASE 1: Syncing SHS Excel with Asset Folders ---")
    # Row 5 is where the headers (Last Name, First Name, etc.) start
    df = pd.read_excel(excel_path, header=5)
    validated_list = []
    
    for index, row in df.iterrows():
        # Clean up input data
        last_name = str(row['Last Name']).strip()
        first_name = str(row['First Name']).strip()
        strand = str(row['STRAND']).strip()
        
        # Combine names to match how we check the folders
        full_name = f"{last_name}, {first_name}"
        
        # 1. Normalize for matching
        norm_surname = normalize_text(last_name)
        norm_strand = normalize_text(strand)
        
        # 2. Navigate to the Strand folder (e.g., G:\Grad\SHS\ABM)
        strand_dir = None
        for p_dir in master_dir.iterdir():
            if p_dir.is_dir() and norm_strand in normalize_text(p_dir.name):
                strand_dir = p_dir
                break
        
        # 3. Find student folder inside that Strand directory
        image_path = None
        if strand_dir:
            for student_folder in strand_dir.iterdir():
                if student_folder.is_dir() and norm_surname in normalize_text(student_folder.name):
                    image_path = next(student_folder.glob('*.[jp][pn]g'), None)
                    break
        
        # 4. Fallback to placeholder if no image
        if not image_path:
            print(f"USING PLACEHOLDER: {full_name}")
            image_path = "no_picture.png" # Ensure this file exists in your folder
            
        validated_list.append({
            "excel_index": index,
            "surname": last_name.upper(),
            "firstname": first_name.title(),
            "course": strand,
            "image_path": image_path
        })
            
    return sorted(validated_list, key=lambda x: x["excel_index"])

# ==========================================
# --- PPTX MANIPULATION ---
# ==========================================
def get_sanitized_image(image_path):
    """Loads image, corrects rotation via EXIF, and converts to stream."""
    with Image.open(image_path) as img:
        img = ImageOps.exif_transpose(img) 
        rgb_img = img.convert('RGB')
        img_stream = io.BytesIO()
        rgb_img.save(img_stream, format="JPEG")
        img_stream.seek(0)
        return img_stream

def duplicate_slide(prs, index):
    """Safely duplicates a slide without corrupting internal relationships."""
    source_slide = prs.slides[index]
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    
    for shp in list(new_slide.shapes):
        shp.element.getparent().remove(shp.element)
        
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(new_el)
            
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        new_slide.part.relate_to(rel.target_part, rel.reltype)
            
    return new_slide

def replace_text_preserve_format(shape, new_text):
    """Injects text into a shape while maintaining font styles."""
    if not shape.has_text_frame:
        return
    
    first_para = shape.text_frame.paragraphs[0]
    first_run = first_para.runs[0] if first_para.runs else first_para.add_run()
    
    saved_font_name = first_run.font.name
    saved_font_size = first_run.font.size
    saved_font_bold = first_run.font.bold
    saved_font_italic = first_run.font.italic
    
    try:
        saved_font_color = first_run.font.color.rgb if first_run.font.color.type else None
    except AttributeError:
        saved_font_color = None

    shape.text_frame.clear()
    
    new_para = shape.text_frame.paragraphs[0]
    new_run = new_para.add_run()
    new_run.text = new_text
    
    new_run.font.name = saved_font_name
    new_run.font.size = saved_font_size
    new_run.font.bold = saved_font_bold
    new_run.font.italic = saved_font_italic
    if saved_font_color:
        new_run.font.color.rgb = saved_font_color

def update_slide(slide, student):
    """Populates a clean slide template with student data."""
    print(f"  -> Assembling: {student['surname']}, {student['firstname']}")
    
    # 1. Update text boxes
    for shape in slide.shapes:
        if shape.name == TARGET_SURNAME_SHAPE:
            replace_text_preserve_format(shape, student["surname"])
        elif shape.name == TARGET_FIRSTNAME_SHAPE:
            replace_text_preserve_format(shape, student["firstname"])
        elif shape.name == TARGET_COURSE_SHAPE:
            replace_text_preserve_format(shape, student["course"])

    # 2. Add Portrait
    clean_image_stream = get_sanitized_image(student["image_path"])
    slide.shapes.add_picture(
        clean_image_stream, 
        IMG_LEFT, 
        IMG_TOP, 
        width=IMG_WIDTH, 
        height=IMG_HEIGHT
    )

# ==========================================
# --- MAIN EXECUTION ---
# ==========================================
def main():
    print("=== COMMENCEMENT AUTOMATION ENGINE STARTED ===")
    
    students = get_validated_students(EXCEL_PATH, MASTER_STUDENT_DIR)
    if not students:
        print("[ERROR] No valid student records found. Aborting.")
        return

    print(f"\n--- PHASE 2: Slide Manufacturing ({len(students)} target slides) ---")
    try:
        prs = Presentation(TEMPLATE_PPTX)
    except Exception as e:
        print(f"[ERROR] Could not open template PPTX: {e}")
        return

    # Generate slides
    for count, student in enumerate(students, start=1):
        print(f"[{count}/{len(students)}]", end=" ")
        new_slide = duplicate_slide(prs, TEMPLATE_SLIDE_INDEX)
        update_slide(new_slide, student)

    # Cleanup: Delete the blank master template slide at the start
    print("\n--- PHASE 3: Cleanup & Output ---")
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

    # Save
    prs.save(OUTPUT_PPTX)
    print(f"\n[SUCCESS] Saved successfully as {OUTPUT_PPTX}")

if __name__ == "__main__":
    main()