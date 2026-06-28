import os
import copy
import io
import re
import cv2
import pandas as pd
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm
from PIL import Image, ImageOps
import unicodedata



# ==========================================
# --- CONFIGURATION & PATHS ---
# ==========================================
EXCEL_PATH = "Graduating_college.xlsx"
MASTER_STUDENT_DIR = Path(r"G:\Grad\College")
TEMPLATE_PPTX = "Commencement Exercises - PPT Template.pptx.pptx"
OUTPUT_PPTX = "Final_Graduation_College.pptx"
TEMPLATE_SLIDE_INDEX = 0
PLACEHOLDER_IMAGE = "no_picture.png"

# ==========================================
# --- SHAPE TARGETS & POSITIONS ---
# ==========================================
TARGET_SURNAME_SHAPE = "TextBox 6"
TARGET_FIRSTNAME_SHAPE = "TextBox 7"
TARGET_COURSE_SHAPE = "TextBox 8"

# Based on PowerPoint 'Format Picture' coordinates
IMG_LEFT = Cm(29.79)   
IMG_TOP = Cm(2.58)     
IMG_WIDTH = Cm(17.27)  
IMG_HEIGHT = Cm(23.43) 

# ==========================================
# --- DATA PROCESSING ---
# ==========================================
def normalize_text(text):
    """Converts Ñ to N and removes special characters for matching."""
    # Convert Ñ to N
    text = unicodedata.normalize('NFKD', text).encode('ASCII', 'ignore').decode('utf-8')
    # Remove underscores and punctuation, convert to lowercase
    return re.sub(r'[^\w\s]', ' ', text).lower()
def parse_name_for_slide(full_name: str):
    """
    Splits the Excel name into Surname and Firstname safely.
    Ideal format in Excel: "De Guzman, Juan"
    """
    if ',' in full_name:
        parts = full_name.split(',', 1)
        return parts[0].strip().upper(), parts[1].strip().title()
    
    # Fallback if no comma is used (Smart PH Name detection)
    words = full_name.split()
    if len(words) == 1:
        return words[0].upper(), ""
        
    prefixes = ['de', 'del', 'de la', 'dela', 'san', 'santa', 'sto', 'sto.']
    
    if len(words) >= 2 and words[-2].lower() in prefixes:
        surname = " ".join(words[-2:])
        first = " ".join(words[:-2])
    elif len(words) >= 3 and " ".join(words[-3:-1]).lower() in prefixes:
        surname = " ".join(words[-3:])
        first = " ".join(words[:-3])
    else:
        surname = words[-1]
        first = " ".join(words[:-1])
        
    return surname.upper(), first.title()

def get_validated_students(excel_path, master_dir):
    print("--- PHASE 1: Syncing Excel with Asset Folders ---")
    df = pd.read_excel(excel_path, header=3)
    validated_list = []
    
    for index, row in df.iterrows():
        excel_name = str(row['STUDENT NAME']).strip()
        program = str(row['PROGRAM']).strip()
        
        # 1. Isolate Surname
        surname_part = excel_name.split(',')[0].strip() if ',' in excel_name else excel_name.split()[-1].strip()
        norm_surname = normalize_text(surname_part)
        norm_prog = normalize_text(program)
        
        # 2. Find folder matching the Program
        program_dir = None
        for p_dir in master_dir.iterdir():
            if p_dir.is_dir() and norm_prog in normalize_text(p_dir.name):
                program_dir = p_dir
                break
        
        # 3. Find the student folder inside that Program directory
        image_path = None
        if program_dir:
            for student_folder in program_dir.iterdir():
                if student_folder.is_dir() and norm_surname in normalize_text(student_folder.name):
                    # Found the folder! Now look for the image
                    image_path = next(student_folder.glob('*.[jp][pn]g'), None)
                    break
        
        # 4. Final Validation: If no image found, use placeholder
        if not image_path:
            print(f"USING PLACEHOLDER: {excel_name}")
            image_path = PLACEHOLDER_IMAGE
            
        surname, firstname = parse_name_for_slide(excel_name)
        validated_list.append({
            "excel_index": index,
            "surname": surname,
            "firstname": firstname,
            "course": program,
            "image_path": image_path
        })
            
    return sorted(validated_list, key=lambda x: x["excel_index"])

def find_best_face(gray_img, face_cascade, img_h):
    """Attempts multiple detection passes from strict to loose, filtering out obvious false positives."""
    # List of strategies: (scaleFactor, minNeighbors, minSize)
    strategies = [
        (1.1, 6, (100, 100)), # Pass 1: Very Strict (Catches perfect, close faces)
        (1.05, 4, (80, 80)),  # Pass 2: Medium (Catches slightly further faces)
        (1.05, 3, (50, 50))   # Pass 3: Loose (Catches far faces, but risky)
    ]
    
    for scale, neighbors, min_s in strategies:
        faces = face_cascade.detectMultiScale(gray_img, scaleFactor=scale, minNeighbors=neighbors, minSize=min_s)
        
        valid_faces = []
        for (x, y, w, h) in faces:
            # SANITY CHECK: Reject any "face" that is in the bottom 30% of the photo
            # This prevents the AI from mistaking a hand, belt, or diploma for a face.
            if y < (img_h * 0.70):
                valid_faces.append((x, y, w, h))
        
        if valid_faces:
            # Return the largest valid face found during this pass
            return max(valid_faces, key=lambda f: f[2] * f[3])
            
    return None # No face found across all passes

# ==========================================
# --- PPTX MANIPULATION ---
# ==========================================
def get_sanitized_image(image_path):
    """Loads image, auto-crops to shoulder/chest height using OpenCV multi-pass, and converts to stream."""
    with Image.open(image_path) as img:
        img = ImageOps.exif_transpose(img)
        img_w, img_h = img.size
        
        rgb_img = img.convert('RGB')
        open_cv_image = np.array(rgb_img)
        gray_img = cv2.cvtColor(open_cv_image, cv2.COLOR_RGB2GRAY)
        
        face_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + 'haarcascade_frontalface_default.xml')
        
        # USE THE NEW MULTI-PASS FUNCTION HERE
        best_face = find_best_face(gray_img, face_cascade, img_h)
        
        target_ratio = 17.27 / 23.43
        
        if best_face is not None:
            # Unpack the coordinates of the validated face
            fx, fy, fw, fh = best_face
            
            # --- DYNAMIC PADDING MATH ---
            top_padding = int(fh * 0.7)     
            bottom_padding = int(fh * 1.8)  # Kept at 1.8 to hide the hands
            
            crop_top = max(0, fy - top_padding)
            crop_bottom = min(img_h, fy + fh + bottom_padding)
            crop_height = crop_bottom - crop_top
            
            crop_width = int(crop_height * target_ratio)
            face_center_x = fx + (fw // 2)
            crop_left = face_center_x - (crop_width // 2)
            crop_right = crop_left + crop_width
            
            if crop_left < 0:
                crop_right -= crop_left
                crop_left = 0
            if crop_right > img_w:
                crop_left -= (crop_right - img_w)
                crop_right = img_w
                
            final_img = rgb_img.crop((crop_left, crop_top, crop_right, crop_bottom))
        else:
            # Fallback (Center Crop)
            print(f"  [!] Face detection completely failed for {image_path}. Using center crop.")
            if img_w / img_h > target_ratio:
                new_w = int(img_h * target_ratio)
                left = (img_w - new_w) // 2
                final_img = rgb_img.crop((left, 0, left + new_w, img_h))
            else:
                new_h = int(img_w / target_ratio)
                top = (img_h - new_h) // 2
                final_img = rgb_img.crop((0, top, img_w, top + new_h))

        img_stream = io.BytesIO()
        final_img.save(img_stream, format="JPEG", quality=95)
        img_stream.seek(0)
        return img_stream

def duplicate_slide(prs, index):
    """Safely duplicates a slide without corrupting internal relationships."""
    source_slide = prs.slides[index]
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    
    for shp in list(new_slide.shapes):
        shp.element.getparent().remove(shp.element)
        
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(new_el)
            
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        new_slide.part.relate_to(rel.target_part, rel.reltype)
            
    return new_slide

def replace_text_preserve_format(shape, new_text):
    """Injects text into a shape while maintaining font styles."""
    if not shape.has_text_frame:
        return
    
    first_para = shape.text_frame.paragraphs[0]
    first_run = first_para.runs[0] if first_para.runs else first_para.add_run()
    
    saved_font_name = first_run.font.name
    saved_font_size = first_run.font.size
    saved_font_bold = first_run.font.bold
    saved_font_italic = first_run.font.italic
    
    try:
        saved_font_color = first_run.font.color.rgb if first_run.font.color.type else None
    except AttributeError:
        saved_font_color = None

    shape.text_frame.clear()
    
    new_para = shape.text_frame.paragraphs[0]
    new_run = new_para.add_run()
    new_run.text = new_text
    
    new_run.font.name = saved_font_name
    new_run.font.size = saved_font_size
    new_run.font.bold = saved_font_bold
    new_run.font.italic = saved_font_italic
    if saved_font_color:
        new_run.font.color.rgb = saved_font_color

def update_slide(slide, student):
    """Populates a clean slide template with student data."""
    print(f"  -> Assembling: {student['surname']}, {student['firstname']}")
    
    # 1. Update text boxes
    for shape in slide.shapes:
        if shape.name == TARGET_SURNAME_SHAPE:
            replace_text_preserve_format(shape, student["surname"])
        elif shape.name == TARGET_FIRSTNAME_SHAPE:
            replace_text_preserve_format(shape, student["firstname"])
        elif shape.name == TARGET_COURSE_SHAPE:
            replace_text_preserve_format(shape, student["course"])

    # 2. Add Portrait
    clean_image_stream = get_sanitized_image(student["image_path"])
    slide.shapes.add_picture(
        clean_image_stream, 
        IMG_LEFT, 
        IMG_TOP, 
        width=IMG_WIDTH, 
        height=IMG_HEIGHT
    )

# ==========================================
# --- MAIN EXECUTION ---
# ==========================================
def main():
    print("=== COMMENCEMENT AUTOMATION ENGINE STARTED ===")
    
    students = get_validated_students(EXCEL_PATH, MASTER_STUDENT_DIR)
    if not students:
        print("[ERROR] No valid student records found. Aborting.")
        return

    print(f"\n--- PHASE 2: Slide Manufacturing ({len(students)} target slides) ---")
    try:
        prs = Presentation(TEMPLATE_PPTX)
    except Exception as e:
        print(f"[ERROR] Could not open template PPTX: {e}")
        return

    # Generate slides
    for count, student in enumerate(students, start=1):
        print(f"[{count}/{len(students)}]", end=" ")
        new_slide = duplicate_slide(prs, TEMPLATE_SLIDE_INDEX)
        update_slide(new_slide, student)

    # Cleanup: Delete the blank master template slide at the start
    print("\n--- PHASE 3: Cleanup & Output ---")
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

    # Save
    prs.save(OUTPUT_PPTX)
    print(f"\n[SUCCESS] Saved successfully as {OUTPUT_PPTX}")

if __name__ == "__main__":
    main()