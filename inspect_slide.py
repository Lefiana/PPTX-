from pptx import Presentation
from pptx.util import Pt
import os

PPTX_FILE = "Commencement Exercises - PPT Template.pptx.pptx"   # Make sure this matches your file name
SLIDE_INDEX = 0                                            # 0 = first slide

def inspect_slide(pptx_path: str, slide_index: int = 0) -> None:
    if not os.path.exists(pptx_path):
        print(f"[ERROR] File not found: {pptx_path}")
        return

    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]

    print("=" * 70)
    print(f"  Slide {slide_index + 1} — Shape Inventory")
    print("=" * 70)
    print(f"{'#':<4} {'Shape Name':<30} {'Type':<20} {'Has Text?':<10} {'Current Text (first 60 chars)'}")
    print("-" * 70)

    for i, shape in enumerate(slide.shapes):
        shape_type = str(shape.shape_type)
        has_text   = shape.has_text_frame

        # Truncate preview text for readability
        text_preview = ""
        if has_text and shape.text_frame.text.strip():
            text_preview = shape.text_frame.text.strip().replace("\n", " ")[:60]

        # Flag placeholders specially
        is_placeholder = shape.is_placeholder
        ph_idx = shape.placeholder_format.idx if is_placeholder else "—"

        print(f"{i:<4} {shape.name:<30} {shape_type:<20} {str(has_text):<10} {text_preview}")

        if is_placeholder:
            print(f"    └─ PLACEHOLDER  idx={ph_idx}")

        # Detailed run-level formatting for text shapes
        if has_text:
            for p_num, para in enumerate(shape.text_frame.paragraphs):
                for r_num, run in enumerate(para.runs):
                    font = run.font
                    # Safely extract font color if it exists
                    color_val = 'inherited'
                    try:
                        if font.color and font.color.type:
                            color_val = font.color.rgb
                    except AttributeError:
                        pass
                        
                    print(
                        f"    └─ Para[{p_num}] Run[{r_num}] "
                        f"bold={font.bold} "
                        f"size={int(font.size.pt) if font.size else 'inherited'}pt "
                        f"color={color_val} "
                        f'text="{run.text[:40]}"'
                    )

    print("=" * 70)

if __name__ == "__main__":
    inspect_slide(PPTX_FILE, SLIDE_INDEX)