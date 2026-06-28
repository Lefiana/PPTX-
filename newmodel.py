import os
import copy
import io
import re
import cv2
import pandas as pd
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.text import MSO_AUTO_SIZE
from PIL import Image, ImageOps
import unicodedata

# ==========================================
# --- CONFIGURATION & PATHS ---
# ==========================================
EXCEL_PATH = "Graduating_college.xlsx"
MASTER_STUDENT_DIR = Path(r"G:\Grad\College")
TEMPLATE_PPTX = "Commencement Exercises - PPT Template.pptx.pptx"
OUTPUT_PPTX = "Final_Graduation_College.pptx"
TEMPLATE_SLIDE_INDEX = 0
PLACEHOLDER_IMAGE = "no_picture.png"

YUNET_MODEL_PATH = "face_detection_yunet_2023mar.onnx"

# ==========================================
# --- SHAPE TARGETS & POSITIONS ---
# ==========================================
TARGET_SURNAME_SHAPE = "TextBox 6"
TARGET_FIRSTNAME_SHAPE = "TextBox 7"
TARGET_COURSE_SHAPE = "TextBox 8"

IMG_LEFT = Cm(29.79)   
IMG_TOP = Cm(2.58)     
IMG_WIDTH = Cm(17.27)  
IMG_HEIGHT = Cm(23.43) 

# ==========================================
# --- DETECTION TRACKER ---
# ==========================================
class DetectionTracker:
    def __init__(self):
        self.total_processed = 0
        self.yunet_success = 0
        self.center_fallback = 0
        self.placeholder_used = 0
    
    def record_yunet_success(self):
        self.total_processed += 1
        self.yunet_success += 1
    
    def record_center_fallback(self):
        self.total_processed += 1
        self.center_fallback += 1
    
    def record_placeholder(self):
        self.total_processed += 1
        self.placeholder_used += 1
    
    def report(self):
        print("\n" + "=" * 60)
        print("--- DETECTION TRACKER REPORT ---")
        print(f"Total Images Processed : {self.total_processed}")
        print(f"YuNet Face Framed      : {self.yunet_success}")
        print(f"Center Crop Fallback   : {self.center_fallback}")
        print(f"Placeholder Images     : {self.placeholder_used}")
        real_photos = self.total_processed - self.placeholder_used
        if real_photos > 0:
            print(f"YuNet Success Rate     : {self.yunet_success / real_photos * 100:.1f}%")
        print("=" * 60)

tracker = DetectionTracker()

# ==========================================
# --- DATA PROCESSING ---
# ==========================================
def normalize_text(text):
    text = unicodedata.normalize('NFKD', text).encode('ASCII', 'ignore').decode('utf-8')
    return re.sub(r'[^\w\s]', ' ', text).lower()

def parse_name_for_slide(full_name: str):
    if ',' in full_name:
        parts = full_name.split(',', 1)
        return parts[0].strip().upper(), parts[1].strip().title()
    
    words = full_name.split()
    if len(words) == 1:
        return words[0].upper(), ""
        
    prefixes = ['de', 'del', 'de la', 'dela', 'san', 'santa', 'sto', 'sto.']
    
    if len(words) >= 2 and words[-2].lower() in prefixes:
        surname = " ".join(words[-2:])
        first = " ".join(words[:-2])
    elif len(words) >= 3 and " ".join(words[-3:-1]).lower() in prefixes:
        surname = " ".join(words[-3:])
        first = " ".join(words[:-3])
    else:
        surname = words[-1]
        first = " ".join(words[:-1])
        
    return surname.upper(), first.title()

def get_validated_students(excel_path, master_dir):
    print("--- PHASE 1: Syncing Excel with Asset Folders ---")
    df = pd.read_excel(excel_path, header=3)
    validated_list = []
    
    for index, row in df.iterrows():
        excel_name = str(row['STUDENT NAME']).strip()
        program = str(row['PROGRAM']).strip()
        
        surname_part = excel_name.split(',')[0].strip() if ',' in excel_name else excel_name.split()[-1].strip()
        norm_surname = normalize_text(surname_part)
        norm_prog = normalize_text(program)
        
        program_dir = None
        for p_dir in master_dir.iterdir():
            if p_dir.is_dir() and norm_prog in normalize_text(p_dir.name):
                program_dir = p_dir
                break
        
        image_path = None
        if program_dir:
            for student_folder in program_dir.iterdir():
                if student_folder.is_dir() and norm_surname in normalize_text(student_folder.name):
                    image_path = next(student_folder.glob('*.[jp][pn]g'), None)
                    break
        
        if not image_path:
            print(f"USING PLACEHOLDER: {excel_name}")
            image_path = PLACEHOLDER_IMAGE
            
        surname, firstname = parse_name_for_slide(excel_name)
        validated_list.append({
            "excel_index": index,
            "surname": surname,
            "firstname": firstname,
            "course": program,
            "image_path": str(image_path)
        })
            
    return sorted(validated_list, key=lambda x: x["excel_index"])

# ==========================================
# --- YUNET FACE DETECTION ---
# ==========================================
def get_yunet_detector(img_w, img_h):
    if not os.path.exists(YUNET_MODEL_PATH):
        raise FileNotFoundError(
            f"YuNet model not found: {YUNET_MODEL_PATH}\n"
            f"Download from: https://github.com/opencv/opencv_zoo/raw/main/models/face_detection_yunet/face_detection_yunet_2023mar.onnx"
        )
    
    detector = cv2.FaceDetectorYN_create(
        YUNET_MODEL_PATH,
        "",
        (img_w, img_h),
        0.3,
        0.3,
        5000
    )
    return detector

def find_best_face_yunet(rgb_img, img_w, img_h):
    """Find the best face using YuNet with smart filtering."""
    try:
        detector = get_yunet_detector(img_w, img_h)
        bgr_img = cv2.cvtColor(np.array(rgb_img), cv2.COLOR_RGB2BGR)
        success, faces = detector.detect(bgr_img)
        
        if not success or faces is None or len(faces) == 0:
            return None
        
        valid_faces = []
        for face in faces:
            x, y, w, h, score = face[0], face[1], face[2], face[3], face[14]
            
            # Skip tiny detections
            if w < 100 or h < 100:
                continue
            
            # Skip faces too low in image (hands, legs, etc.)
            face_center_y = y + (h / 2)
            if face_center_y > img_h * 0.65:
                continue
            
            # Skip faces that are too small or too large relative to image
            face_area_ratio = (w * h) / (img_w * img_h)
            if face_area_ratio < 0.005 or face_area_ratio > 0.5:
                continue
            
            # Combined scoring: confidence + position + size
            position_score = 1.0 - (face_center_y / img_h)
            size_score = min(face_area_ratio * 20, 1.0)
            combined_score = (score * 0.5) + (position_score * 0.3) + (size_score * 0.2)
            
            valid_faces.append((int(x), int(y), int(w), int(h), float(score), float(combined_score)))
        
        if not valid_faces:
            return None
        
        best = max(valid_faces, key=lambda f: f[5])
        return best[:4]
        
    except Exception as e:
        print(f"  [YuNet ERROR] {str(e)[:80]}")
        return None

# ==========================================
# --- PPTX MANIPULATION ---
# ==========================================
def get_sanitized_image(image_path):
    filename = os.path.basename(image_path)
    
    with Image.open(image_path) as img:
        img = ImageOps.exif_transpose(img)
        img_w, img_h = img.size
        
        rgb_img = img.convert('RGB')
        target_ratio = 17.27 / 23.43  # ~0.737
        
        is_placeholder = (image_path == PLACEHOLDER_IMAGE or PLACEHOLDER_IMAGE in image_path)
        
        best_face = None
        if not is_placeholder:
            best_face = find_best_face_yunet(rgb_img, img_w, img_h)
        
        if best_face is not None:
            fx, fy, fw, fh = best_face
            
            # --- PORTRAIT FRAMING BASED ON REFERENCE IMAGE ---
            # Reference shows: head + upper torso, eyes at top third, centered face
            # 
            # Key measurements from reference:
            # - Headroom above head: ~15-20% of crop height
            # - Eyes position: ~30% from top of crop (top third rule)
            # - Bottom: just below chest/medallion level
            # - Face horizontally centered
            
            # YuNet gives us face box (fx, fy, fw, fh)
            # We estimate eye position as upper portion of face box
            # Typically eyes are at ~30% from top of face box
            eye_y = fy + int(fh * 0.3)  # Estimated eye level
            face_cx = fx + (fw / 2)
            
            # Portrait composition: eyes at ~30% from top of final crop
            EYE_VERTICAL_POSITION = 0.30
            
            # Calculate crop height: we want eyes at 30% from top
            # So total height = eye_y / 0.30, but also need room below
            # For graduation portrait: show down to chest (~2.5x face height below eyes)
            
            # Height needed above eyes
            height_above_eyes = int(eye_y / EYE_VERTICAL_POSITION)
            # Height needed below eyes (show chest/medallion area)
            height_below_eyes = int(fh * 3.5)  # 2.8x face height below eyes
            
            crop_height = height_above_eyes + height_below_eyes
            crop_height = min(crop_height, img_h)  # Don't exceed image
            
            # Calculate top position
            crop_top = max(0, eye_y - int(crop_height * EYE_VERTICAL_POSITION))
            
            # If we hit the top, recalculate
            if crop_top == 0:
                # Eyes are near top of image, use fixed proportion
                crop_height = min(int(img_h * 0.72), img_h)  # Show 72% of image height
                crop_top = 0
            else:
                # Ensure we don't go below image bottom
                if crop_top + crop_height > img_h:
                    crop_height = img_h - crop_top
            
            crop_bottom = crop_top + crop_height
            
            # Calculate width from height and target ratio
            crop_width = int(crop_height * target_ratio)
            
            # Center horizontally on face center
            crop_left = int(face_cx - (crop_width / 2))
            crop_right = crop_left + crop_width
            
            # Keep within bounds
            if crop_left < 0:
                crop_right -= crop_left
                crop_left = 0
            if crop_right > img_w:
                crop_left -= (crop_right - img_w)
                crop_right = img_w
            
            # Final validation
            if (crop_width < 200 or crop_height < 200 or 
                crop_left < 0 or crop_right > img_w or 
                crop_bottom > img_h or crop_top < 0):
                print(f"  [!] Invalid crop for {filename}. Using center crop.")
                best_face = None
            else:
                final_img = rgb_img.crop((crop_left, crop_top, crop_right, crop_bottom))
                tracker.record_yunet_success()
                print(f"  [OK] Portrait crop: eyes@{EYE_VERTICAL_POSITION:.0%}, {crop_width}x{crop_height}")
        
        # Fallback: center crop
        if best_face is None:
            if is_placeholder:
                tracker.record_placeholder()
            else:
                tracker.record_center_fallback()
                print(f"  [!] Center crop for {filename}.")
            
            if img_w / img_h > target_ratio:
                new_w = int(img_h * target_ratio)
                left = (img_w - new_w) // 2
                final_img = rgb_img.crop((left, 0, left + new_w, img_h))
            else:
                new_h = int(img_w / target_ratio)
                top = (img_h - new_h) // 2
                final_img = rgb_img.crop((0, top, img_w, top + new_h))

        img_stream = io.BytesIO()
        final_img.save(img_stream, format="JPEG", quality=95)
        img_stream.seek(0)
        return img_stream

def duplicate_slide(prs, index):
    source_slide = prs.slides[index]
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    
    for shp in list(new_slide.shapes):
        shp.element.getparent().remove(shp.element)
        
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(new_el)
            
    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        new_slide.part.relate_to(rel.target_part, rel.reltype)
            
    return new_slide

def replace_text_preserve_format(shape, new_text):
    if not shape.has_text_frame:
        return
    
    shape.text_frame.word_wrap = False
    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    first_para = shape.text_frame.paragraphs[0]
    first_run = first_para.runs[0] if first_para.runs else first_para.add_run()
    
    saved_font_name = first_run.font.name
    saved_font_size = first_run.font.size
    saved_font_bold = first_run.font.bold
    saved_font_italic = first_run.font.italic
    
    try:
        saved_font_color = first_run.font.color.rgb if first_run.font.color.type else None
    except AttributeError:
        saved_font_color = None

    shape.text_frame.clear()
    
    new_para = shape.text_frame.paragraphs[0]
    new_run = new_para.add_run()
    new_run.text = new_text
    
    new_run.font.name = saved_font_name
    new_run.font.size = saved_font_size
    new_run.font.bold = saved_font_bold
    new_run.font.italic = saved_font_italic
    if saved_font_color:
        new_run.font.color.rgb = saved_font_color

def update_slide(slide, student):
    print(f"  -> Assembling: {student['surname']}, {student['firstname']}")
    
    for shape in slide.shapes:
        if shape.name == TARGET_SURNAME_SHAPE:
            replace_text_preserve_format(shape, student["surname"])
        elif shape.name == TARGET_FIRSTNAME_SHAPE:
            replace_text_preserve_format(shape, student["firstname"])
        elif shape.name == TARGET_COURSE_SHAPE:
            replace_text_preserve_format(shape, student["course"])

    clean_image_stream = get_sanitized_image(student["image_path"])
    slide.shapes.add_picture(
        clean_image_stream, 
        IMG_LEFT, 
        IMG_TOP, 
        width=IMG_WIDTH, 
        height=IMG_HEIGHT
    )

# ==========================================
# --- MAIN EXECUTION ---
# ==========================================
def main():
    print("=== COMMENCEMENT AUTOMATION ENGINE STARTED ===")
    print(f"OpenCV version: {cv2.__version__}")
    
    if not os.path.exists(YUNET_MODEL_PATH):
        print(f"[ERROR] YuNet model not found at: {YUNET_MODEL_PATH}")
        print("Please download it from:")
        print("https://github.com/opencv/opencv_zoo/raw/main/models/face_detection_yunet/face_detection_yunet_2023mar.onnx")
        return
    
    print(f"YuNet model size: {os.path.getsize(YUNET_MODEL_PATH):,} bytes")
    
    students = get_validated_students(EXCEL_PATH, MASTER_STUDENT_DIR)
    if not students:
        print("[ERROR] No valid student records found. Aborting.")
        return

    print(f"\n--- PHASE 2: Slide Manufacturing ({len(students)} target slides) ---")
    try:
        prs = Presentation(TEMPLATE_PPTX)
    except Exception as e:
        print(f"[ERROR] Could not open template PPTX: {e}")
        return

    for count, student in enumerate(students, start=1):
        print(f"[{count}/{len(students)}]", end=" ")
        new_slide = duplicate_slide(prs, TEMPLATE_SLIDE_INDEX)
        update_slide(new_slide, student)

    print("\n--- PHASE 3: Cleanup & Output ---")
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

    prs.save(OUTPUT_PPTX)
    
    tracker.report()
    
    print(f"\n[SUCCESS] Saved successfully as {OUTPUT_PPTX}")

if __name__ == "__main__":
    main()